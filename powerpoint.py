import streamlit as st
import os
import shutil
import zipfile
import uuid
from pathlib import Path
from pptx import Presentation
from jinja2 import Template

# Constants
TEMPLATES_DIR = Path("templates")
STATIC_DIR = Path("static")
OUTPUT_DIR = Path("output")

# Setup dirs
OUTPUT_DIR.mkdir(exist_ok=True)


def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    slides = []
    for i, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        slides.append({"index": i + 1, "text": text.strip()})
    return slides


def generate_slide_html(slide, total, timer_enabled):
    with open(TEMPLATES_DIR / "slide_template.html") as f:
        template = Template(f.read())
    return template.render(slide=slide, total=total, timer_enabled=timer_enabled)


def generate_index_html(slides, duration_sec, mode):
    with open(TEMPLATES_DIR / "index_template.html") as f:
        template = Template(f.read())
    return template.render(slides=slides, duration=duration_sec, mode=mode)


def generate_manifest(scorm_version, title):
    template_file = "imsmanifest_1.2.xml.j2" if scorm_version == "1.2" else "imsmanifest_2004.xml.j2"
    with open(STATIC_DIR / template_file) as f:
        template = Template(f.read())
    return template.render(title=title)


def create_scorm_package(slides, scorm_version, duration_sec, mode):
    package_id = str(uuid.uuid4())
    package_path = OUTPUT_DIR / f"scorm_{package_id}"
    os.makedirs(package_path / "slides")

    # Copy static files
    shutil.copy(STATIC_DIR / "scorm_api_wrapper.js", package_path)

    # Generate slide files
    for slide in slides:
        html = generate_slide_html(slide, len(slides), timer_enabled=(mode in ["duration", "both"]))
        with open(package_path / "slides" / f"slide_{slide['index']}.html", "w") as f:
            f.write(html)

    # Generate index.html
    index_html = generate_index_html(slides, duration_sec, mode)
    with open(package_path / "index.html", "w") as f:
        f.write(index_html)

    # Manifest
    manifest = generate_manifest(scorm_version, title="Converted PPTX")
    with open(package_path / "imsmanifest.xml", "w") as f:
        f.write(manifest)

    # Zip
    zip_path = str(package_path) + ".zip"
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, _, files in os.walk(package_path):
            for file in files:
                file_path = Path(root) / file
                arcname = file_path.relative_to(package_path)
                zipf.write(file_path, arcname)

    return zip_path


# Streamlit interface
st.title("Convertisseur PowerPoint en SCORM")

uploaded_file = st.file_uploader("Uploader un fichier .pptx", type="pptx")
scorm_version = st.selectbox("Choisir la version SCORM", ["1.2", "2004"])
duration_min = st.number_input("Durée minimale (minutes) pour valider la complétion", min_value=0, value=1)
mode = st.selectbox("Condition de complétion", ["duration", "slides", "both"])

if st.button("Convertir") and uploaded_file:
    with st.spinner("Conversion en cours..."):
        slides = extract_text_from_pptx(uploaded_file)
        zip_path = create_scorm_package(slides, scorm_version, duration_min * 60, mode)
        st.success("Conversion réussie !")
        with open(zip_path, "rb") as f:
            st.download_button("Télécharger le package SCORM", f, file_name=os.path.basename(zip_path))
